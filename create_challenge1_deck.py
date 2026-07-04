"""
create_challenge1_deck.py
Builds: Challenge1_Indian_English_Piper_TTS.pptx  (15 slides, 16:9)

A presentation deck for CTO Hackathon — Challenge 1 (Indian English Voice for Piper TTS).
Every slide uses plain, explainable English that maps 1:1 to vocabulary_glossary.md,
so the presenter can expand any term verbally.

Palette + helpers adapted from ../../Bearpaw/create_atticus_deck.py

Run:
    pip install python-pptx
    python create_challenge1_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = Path(__file__).parent
OUTPUT_PATH = BASE / "Challenge1_Indian_English_Piper_TTS.pptx"

# ── Brand palette ────────────────────────────────────────────────────────────
DARK_BG  = RGBColor(  5,   7,  10)
CARD_BG  = RGBColor( 13,  19,  31)
CARD_BG2 = RGBColor( 15,  23,  42)
BORDER   = RGBColor( 30,  41,  59)
CYAN     = RGBColor(  0, 240, 255)
PURPLE   = RGBColor(112,   0, 255)
VIOLET   = RGBColor(139,  92, 246)
BLUE     = RGBColor( 59, 130, 246)
EMERALD  = RGBColor( 52, 211, 153)
AMBER    = RGBColor(251, 191,  36)
SKY      = RGBColor( 56, 189, 248)
PINK     = RGBColor(236,  72, 153)
RED      = RGBColor(255,  71,  87)
WHITE    = RGBColor(255, 255, 255)
SL200    = RGBColor(226, 232, 240)
SL400    = RGBColor(148, 163, 184)
SL500    = RGBColor(100, 116, 139)
SL600    = RGBColor( 71,  85, 105)

SW = Inches(13.333)
SH = Inches(7.5)

from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ── Core helpers ─────────────────────────────────────────────────────────────
def set_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_rect(slide, l, t, w, h, fill=CARD_BG, line=BORDER, line_w=Pt(0.75), radius=0.08):
    if radius > 0:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_txt(slide, text, l, t, w, h, size=Pt(12), color=SL200, bold=False,
            align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def add_bullets(slide, items, l, t, w, h, size=Pt(12), color=SL400,
                gap=6, bullet="•  ", lead_color=CYAN):
    """items: list of str OR (lead, rest) tuples (lead rendered bold + coloured)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run(); r1.text = bullet + lead
            r1.font.size = size; r1.font.bold = True; r1.font.color.rgb = lead_color
            r2 = p.add_run(); r2.text = rest
            r2.font.size = size; r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = bullet + item
            r.font.size = size; r.font.color.rgb = color
    return tb


def accent_bar(slide):
    bar_h = Inches(0.055)
    add_rect(slide, 0, 0, SW // 2, bar_h, fill=CYAN, line=None, radius=0)
    add_rect(slide, SW // 2, 0, SW // 2, bar_h, fill=PURPLE, line=None, radius=0)
    add_rect(slide, SW // 2 - Inches(0.75), 0, Inches(1.5), bar_h, fill=VIOLET, line=None, radius=0)


def deco_circles(slide):
    for cx, cy, r, col in [
        (Inches(-1),   Inches(7.5), Inches(3.5), CYAN),
        (Inches(13.5), Inches(0),   Inches(3),   PURPLE),
        (Inches(12),   Inches(8),   Inches(2.5), RED),
    ]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
        circle.fill.solid()
        circle.fill.fore_color.rgb = col
        circle.line.fill.background()
        circle.shadow.inherit = False
        sp_pr = circle.fill._xPr
        solid = sp_pr.find(qn('a:solidFill'))
        if solid is not None:
            srgb = solid.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = etree.SubElement(srgb, qn('a:alpha'))
                alpha.set('val', str(int(0.04 * 100000)))


def header(slide, kicker, title, subtitle=None):
    add_txt(slide, kicker, Inches(0.55), Inches(0.30), Inches(11), Inches(0.3),
            size=Pt(11), color=CYAN, bold=True)
    add_txt(slide, title, Inches(0.55), Inches(0.58), Inches(12.2), Inches(0.7),
            size=Pt(27), color=WHITE, bold=True)
    if subtitle:
        add_txt(slide, subtitle, Inches(0.55), Inches(1.20), Inches(12.2), Inches(0.4),
                size=Pt(12.5), color=SL400)
    add_rect(slide, Inches(0.55), Inches(1.62), Inches(12.23), Pt(1.4), fill=BORDER, line=None, radius=0)


def footer(slide, n):
    add_txt(slide, "Indian English Voice for Piper TTS  ·  CTO Hackathon — Challenge 1",
            Inches(0.55), Inches(7.12), Inches(10), Inches(0.3), size=Pt(8), color=SL600)
    add_txt(slide, f"{n:02d} / 15", Inches(11.9), Inches(7.12), Inches(0.9), Inches(0.3),
            size=Pt(8), color=SL600, align=PP_ALIGN.RIGHT)


def feature_card(slide, x, y, w, h, col, title, body, title_size=Pt(14), body_size=Pt(10.5)):
    add_rect(slide, x, y, w, h, fill=CARD_BG2, line=col, line_w=Pt(1), radius=0.05)
    add_rect(slide, x, y, Inches(0.06), h, fill=col, line=None, radius=0)
    add_txt(slide, title, x + Inches(0.24), y + Inches(0.15), w - Inches(0.34), Inches(0.4),
            size=title_size, color=col, bold=True)
    by = y + Inches(0.62)
    if isinstance(body, list):
        add_bullets(slide, body, x + Inches(0.24), by, w - Inches(0.42), h - Inches(0.72),
                    size=body_size, color=SL400, gap=4, lead_color=col)
    else:
        add_txt(slide, body, x + Inches(0.24), by, w - Inches(0.42), h - Inches(0.72),
                size=body_size, color=SL400)


def flow_boxes(slide, y, items, h=Inches(1.0), start=Inches(0.6), total_w=Inches(12.13)):
    """items: (label, color, sublabel) tuples. Lays boxes left→right with → arrows."""
    n = len(items)
    arrow_w = Inches(0.5)
    box_w = (total_w - arrow_w * (n - 1)) // n
    x = start
    for i, (label, col, sub) in enumerate(items):
        add_rect(slide, x, y, box_w, h, fill=CARD_BG2, line=col, line_w=Pt(1.25), radius=0.12)
        add_txt(slide, label, x + Inches(0.05), y + Inches(0.14), box_w - Inches(0.1), Inches(0.5),
                size=Pt(12.5), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        if sub:
            add_txt(slide, sub, x + Inches(0.05), y + Inches(0.56), box_w - Inches(0.1), Inches(0.4),
                    size=Pt(8.5), color=col, align=PP_ALIGN.CENTER)
        if i < n - 1:
            add_txt(slide, "→", x + box_w, y, arrow_w, h, size=Pt(22), color=SL500,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = x + box_w + arrow_w


def callout(slide, text, y, col=EMERALD, l=Inches(0.6), w=Inches(12.13), h=Inches(0.62)):
    add_rect(slide, l, y, w, h, fill=RGBColor(col[0] // 7, col[1] // 7, col[2] // 7),
             line=col, line_w=Pt(1), radius=0.2)
    add_txt(slide, text, l + Inches(0.25), y, w - Inches(0.5), h,
            size=Pt(12), color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    deco_circles(s)
    accent_bar(s)
    return s


# ── Slides ───────────────────────────────────────────────────────────────────
def slide_01_title(prs):
    s = new_slide(prs)
    add_txt(s, "CTO HACKATHON  ·  CHALLENGE 1", Inches(1.5), Inches(1.7), Inches(10.33), Inches(0.4),
            size=Pt(13), color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_txt(s, "Indian English Voice", Inches(1), Inches(2.35), Inches(11.33), Inches(1.0),
            size=Pt(54), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_txt(s, "for Piper TTS", Inches(1), Inches(3.35), Inches(11.33), Inches(0.8),
            size=Pt(40), color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
    add_txt(s, "A complete technical approach — from phonemes to a browser-ready neural voice",
            Inches(1.5), Inches(4.35), Inches(10.33), Inches(0.5),
            size=Pt(15), color=SL400, align=PP_ALIGN.CENTER)

    pills = [("Piper TTS", CYAN), ("eSpeak-ng", VIOLET), ("VITS", EMERALD),
             ("ONNX Runtime Web", AMBER), ("Fine-Tuning", SKY), ("100% Client-Side", PINK)]
    pw, gap = Inches(1.78), Inches(0.16)
    total = len(pills) * pw + (len(pills) - 1) * gap
    x0 = (SW - total) // 2
    for i, (label, col) in enumerate(pills):
        x = x0 + i * (pw + gap)
        sh = add_rect(s, x, Inches(5.25), pw, Inches(0.42),
                      fill=RGBColor(col[0] // 6, col[1] // 6, col[2] // 6),
                      line=col, line_w=Pt(1), radius=0.5)
        tf = sh.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(10); r.font.color.rgb = col; r.font.bold = True
    add_txt(s, "Piper TTS  +  eSpeak-ng  —  open source, on-device",
            Inches(1), Inches(6.55), Inches(11.33), Inches(0.4),
            size=Pt(10), color=SL600, align=PP_ALIGN.CENTER)


def slide_02_problem(prs):
    s = new_slide(prs)
    header(s, "THE CHALLENGE", "Understandable — but not truly Indian",
           "Today's voice is clear, yet it doesn't sound authentically Indian English.")
    feature_card(s, Inches(0.55), Inches(1.9), Inches(6.0), Inches(2.15), AMBER,
                 "The problem today",
                 ["Voice is clear but carries a US / UK accent",
                  "Numbers, names and loanwords sound wrong",
                  "“Lakh / crore”, “Aadhaar”, Indian names misread"])
    feature_card(s, Inches(6.78), Inches(1.9), Inches(6.0), Inches(2.15), EMERALD,
                 "What we want",
                 ["A natural Indian-English voice",
                  "Right sounds, rhythm and intonation",
                  "Clear, human, and locally recognisable"])
    feature_card(s, Inches(0.55), Inches(4.25), Inches(6.0), Inches(2.1), CYAN,
                 "The hard constraint",
                 ["Must run 100% in the browser — no server",
                  "Low latency, works offline",
                  "Small enough to download quickly"])
    feature_card(s, Inches(6.78), Inches(4.25), Inches(6.0), Inches(2.1), VIOLET,
                 "The toolkit",
                 [("Piper TTS ", "— the neural voice engine"),
                  ("eSpeak-ng ", "— turns text into phonemes"),
                  ("Both open source, both run on-device", "")])
    footer(s, 2)


def slide_03_bigidea(prs):
    s = new_slide(prs)
    header(s, "THE CORE IDEA", "Speech happens in two stages",
           "Understand this split and every design decision follows.")
    feature_card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(3.4), CYAN,
                 "1 · eSpeak-ng  —  “the script”",
                 ["Decides WHICH sounds to make (the phonemes)",
                  "Rule-based, tiny, and fast",
                  "Turns letters into pronunciation symbols",
                  "Example: “crore” → /kroːr/"])
    feature_card(s, Inches(6.78), Inches(1.95), Inches(6.0), Inches(3.4), VIOLET,
                 "2 · VITS neural model  —  “the performance”",
                 ["Decides HOW those sounds are voiced",
                  "Controls accent, pitch, rhythm, timbre",
                  "Learns the Indian voice from real recordings",
                  "This is where the accent truly lives"])
    callout(s, "Indian English differs in BOTH stages — so our solution addresses both.  "
               "Think of an actor (VITS) performing a script (eSpeak-ng).", Inches(5.55), col=EMERALD)
    footer(s, 3)


def slide_04_pipeline(prs):
    s = new_slide(prs)
    header(s, "TASK 1 · PIPELINE ANALYSIS", "From text to voice, step by step")
    flow_boxes(s, Inches(2.15), [
        ("Input Text", SL400, "“150 crore”"),
        ("eSpeak-ng", CYAN, "text → phonemes"),
        ("Phoneme → ID", SKY, "sounds → numbers"),
        ("VITS Model", VIOLET, "numbers → audio"),
        ("Waveform", EMERALD, "you hear it \U0001F50A"),
    ])
    feature_card(s, Inches(0.55), Inches(3.6), Inches(6.0), Inches(1.75), CYAN,
                 "Training (build the voice)",
                 ["Feed phonemes + real audio to VITS",
                  "The model learns to map sounds → voice"])
    feature_card(s, Inches(6.78), Inches(3.6), Inches(6.0), Inches(1.75), VIOLET,
                 "Inference (use the voice)",
                 ["Same steps run live in the browser",
                  "eSpeak-ng (WASM) → VITS (ONNX) → speaker"])
    callout(s, "The golden rule: phonemization at inference must EXACTLY match training — "
               "same voice, same dictionary — or the model hears sounds it never learned.",
            Inches(5.55), col=AMBER)
    footer(s, 4)


def slide_05_levers(prs):
    s = new_slide(prs)
    header(s, "TASK 1 · WHAT DRIVES QUALITY", "Three levers of quality")
    feature_card(s, Inches(0.55), Inches(1.95), Inches(3.95), Inches(4.15), CYAN,
                 "Pronunciation",
                 [("Owned by eSpeak-ng", ""),
                  "The lexicon + letter-to-sound rules pick the phonemes",
                  "Wrong phoneme → the model faithfully says the wrong sound",
                  "Fix: dictionary entries for Indian words"])
    feature_card(s, Inches(4.69), Inches(1.95), Inches(3.95), Inches(4.15), VIOLET,
                 "Accent",
                 [("Owned by VITS", ""),
                  "Duration + pitch predictors learn Indian rhythm & intonation",
                  "Syllable-timed cadence, retroflex ‘t’/‘d’",
                  "Fix: fine-tune on Indian speech"])
    feature_card(s, Inches(8.83), Inches(1.95), Inches(3.95), Inches(4.15), EMERALD,
                 "Naturalness",
                 [("Owned by VITS", ""),
                  "HiFi-GAN vocoder turns features into a real waveform",
                  "Normalizing flows add human variation",
                  "Higher sample rate = richer audio"])
    footer(s, 5)


def slide_06_eninsight(prs):
    s = new_slide(prs)
    header(s, "THE KEY INSIGHT", "There is no “Indian English” in eSpeak-ng",
           "A technical judge will check this — so we turn it into a strength.")
    feature_card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(3.4), AMBER,
                 "What eSpeak-ng actually ships",
                 ["en-us  —  American", "en  —  British (default)",
                  "en-029  —  Caribbean (NOT Indian)",
                  "a few UK regionals",
                  "✖  No ‘en-in’.  ‘-v en-in’ simply fails."])
    feature_card(s, Inches(6.78), Inches(1.95), Inches(6.0), Inches(3.4), EMERALD,
                 "So how do we get the accent?",
                 ["We don't rely on eSpeak for it",
                  "Phonemize with standard en-us",
                  "Let the NEURAL model learn the Indian accent",
                  "…from real Indian-English recordings"])
    callout(s, "Naming this gap — and solving it correctly — is exactly the “deep understanding” "
               "the challenge rewards.", Inches(5.55), col=CYAN)
    footer(s, 6)


def slide_07_twopaths(prs):
    s = new_slide(prs)
    header(s, "OUR APPROACH", "Two ways to get the accent")
    feature_card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(3.4), EMERALD,
                 "Path A  —  Acoustic  (Recommended)",
                 ["Phonemize with the stock en-us voice",
                  "The neural model learns the accent from data",
                  "Simple, robust, no eSpeak surgery",
                  "Nothing extra to ship to the browser"])
    feature_card(s, Inches(6.78), Inches(1.95), Inches(6.0), Inches(3.4), AMBER,
                 "Path B  —  Custom eSpeak  (Advanced)",
                 ["Author Indian rules so eSpeak emits retroflex sounds",
                  "Add a dictionary for names & loanwords",
                  "Higher fidelity, but more engineering",
                  "Must ship the custom data in the browser too"])
    callout(s, "Recommendation: start with Path A. Layer in a custom loanword dictionary "
               "(high value, low risk) before full retroflex rules.", Inches(5.55), col=CYAN)
    footer(s, 7)


def slide_08_data(prs):
    s = new_slide(prs)
    header(s, "TASK 2 · VOICE DEVELOPMENT STRATEGY", "Data: the ceiling on quality",
           "No amount of training fixes noisy or inconsistent audio.")
    feature_card(s, Inches(0.55), Inches(1.9), Inches(6.0), Inches(2.5), CYAN,
                 "Primary  —  IndicTTS (IIT Madras)",
                 ["Studio-recorded, single speaker, very clean",
                  "Phonetically balanced for speech synthesis",
                  "Best for a natural single-speaker voice",
                  "Available for research & development"])
    feature_card(s, Inches(6.78), Inches(1.9), Inches(6.0), Inches(2.5), VIOLET,
                 "Alternative  —  Common Voice (India)",
                 ["Many speakers, realistic urban accents",
                  "Permissive CC0 license — safe for commercial use",
                  "Noisier → needs heavier cleaning",
                  "Best when speaker variety matters most"])
    add_txt(s, "Data preparation pipeline", Inches(0.55), Inches(4.55), Inches(8), Inches(0.35),
            size=Pt(13), color=WHITE, bold=True)
    flow_boxes(s, Inches(4.95), [
        ("Resample", SKY, "16 / 22.05 kHz"),
        ("Normalize", EMERALD, "even loudness"),
        ("Trim silence", AMBER, "clean edges"),
        ("Normalize text", CYAN, "1,50,000 → ‘one lakh…’"),
    ], h=Inches(0.95))
    footer(s, 8)


def slide_09_finetune(prs):
    s = new_slide(prs)
    header(s, "TASK 2 · ADAPTATION STRATEGY", "Fine-tune — don't start from scratch")
    feature_card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(2.7), RED,
                 "✖  Training from scratch",
                 ["40–100+ hours of studio audio",
                  "Weeks of training on many GPUs",
                  "Must relearn how to make sound at all",
                  "Overkill for one accent"])
    feature_card(s, Inches(6.78), Inches(1.95), Inches(6.0), Inches(2.7), EMERALD,
                 "✔  Fine-tuning (our choice)",
                 ["Only 3–5 hours of clean Indian audio",
                  "4–6 hours on a single GPU (~10–20k steps)",
                  "Starts from en_US-lessac-medium",
                  "Only shifts accent, pitch & rhythm"])
    callout(s, "Why it works: the vocoder already makes clean audio — we just teach it the "
               "Indian accent. Fine-tuning is a gentle nudge, not a rebuild.", Inches(4.95), col=CYAN)
    footer(s, 9)


def slide_10_hyperparams(prs):
    s = new_slide(prs)
    header(s, "TASK 3 · TRAINING SETTINGS", "The knobs that matter — and why")
    rows = [
        ("Learning rate", "2 × 10⁻⁴", "How big each step is — too high wrecks the voice"),
        ("Batch size", "16 – 32", "Samples per step — balances speed vs memory"),
        ("LR decay", "0.9998 / epoch", "Slows down near the finish to avoid overshoot"),
        ("Grad clip", "1.0", "Caps updates — keeps VITS training stable"),
        ("Warmup", "~1,000 steps", "Eases in gently, protecting the pre-trained model"),
        ("Fine-tune steps", "10k – 20k", "Enough to adapt, not enough to overfit"),
    ]
    y = Inches(1.95)
    rh = Inches(0.72)
    for i, (name, val, why) in enumerate(rows):
        add_rect(s, Inches(0.55), y, Inches(12.23), rh - Inches(0.1),
                 fill=CARD_BG2 if i % 2 == 0 else CARD_BG, line=BORDER, line_w=Pt(0.5), radius=0.06)
        add_txt(s, name, Inches(0.75), y, Inches(3.1), rh - Inches(0.1),
                size=Pt(13), color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_txt(s, val, Inches(3.9), y, Inches(2.4), rh - Inches(0.1),
                size=Pt(13), color=CYAN, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_txt(s, why, Inches(6.4), y, Inches(6.2), rh - Inches(0.1),
                size=Pt(11), color=SL400, anchor=MSO_ANCHOR.MIDDLE)
        y = y + rh
    callout(s, "The theme: a gentle nudge that PRESERVES the strong pre-trained model.",
            Inches(6.5), col=EMERALD)
    footer(s, 10)


def slide_11_eval(prs):
    s = new_slide(prs)
    header(s, "TASK 3 · MEASURING SUCCESS", "How we know the voice is good")
    feature_card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(3.4), CYAN,
                 "Objective  (automatic)",
                 [("WER < 5%  ", "— clarity: transcribe it back, count errors"),
                  ("MCD < 3.0  ", "— similarity to the real voice"),
                  "Cheap, repeatable — great as a training gate"])
    feature_card(s, Inches(6.78), Inches(1.95), Inches(6.0), Inches(3.4), VIOLET,
                 "Subjective  (humans)",
                 [("MOS 1–5  ", "— people rate it out of 5"),
                  "Intelligibility · Naturalness · Accent authenticity",
                  ("A/B test  ", "— base model vs fine-tuned model")])
    callout(s, "Numbers gate the training; humans have the final word — a voice can score well "
               "yet still sound foreign, and only listeners catch that.", Inches(5.55), col=AMBER)
    footer(s, 11)


def slide_12_deploy(prs):
    s = new_slide(prs)
    header(s, "TASK 4 · DEPLOYMENT", "Runs 100% in the browser")
    flow_boxes(s, Inches(2.0), [
        ("INT8 ONNX", AMBER, "~15–18 MB"),
        ("ONNX Runtime Web", CYAN, "WASM + SIMD"),
        ("Web Worker", VIOLET, "off the UI thread"),
        ("Web Audio", EMERALD, "< 150 ms to sound"),
    ])
    feature_card(s, Inches(0.55), Inches(3.45), Inches(3.95), Inches(1.95), AMBER,
                 "Small",
                 ["INT8 makes the model ~70% smaller",
                  "Downloads fast, low memory"])
    feature_card(s, Inches(4.69), Inches(3.45), Inches(3.95), Inches(1.95), CYAN,
                 "Fast",
                 ["SIMD + threads accelerate it",
                  "Faster than real time (RTF < 1)"])
    feature_card(s, Inches(8.83), Inches(3.45), Inches(3.95), Inches(1.95), EMERALD,
                 "Smooth & offline",
                 ["Web Worker keeps the UI responsive",
                  "Cached in IndexedDB after first load"])
    callout(s, "No server, no round-trip, no privacy concern — the whole voice lives in the browser tab.",
            Inches(5.65), col=VIOLET)
    footer(s, 12)


def slide_13_tradeoffs(prs):
    s = new_slide(prs)
    header(s, "TRADE-OFFS & CHALLENGES", "What we watch for")
    feature_card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(2.0), CYAN,
                 "Quality vs latency",
                 ["Ship Medium (22 kHz) by default",
                  "Auto-fallback to Low (16 kHz) on weak devices"])
    feature_card(s, Inches(6.78), Inches(1.95), Inches(6.0), Inches(2.0), VIOLET,
                 "Accent drift",
                 ["Model can slip back toward US sounds",
                  "Fix: a balanced training set + A/B checks"])
    feature_card(s, Inches(0.55), Inches(4.1), Inches(6.0), Inches(2.0), AMBER,
                 "Code-mixed words",
                 ["“Aadhaar”, “crore”, names like “Goutam”",
                  "Fix: a custom loanword dictionary"])
    feature_card(s, Inches(6.78), Inches(4.1), Inches(6.0), Inches(2.0), RED,
                 "Train = Inference  (the subtle one)",
                 ["Browser must use the SAME eSpeak + dictionary",
                  "Mismatch → garbled speech"])
    footer(s, 13)


def slide_14_recommendation(prs):
    s = new_slide(prs)
    header(s, "FINAL RECOMMENDATION", "The plan, in one view")
    rows = [
        ("Strategy", "Fine-tune a pre-trained model (not from scratch)", EMERALD),
        ("Base model", "en_US-lessac-medium — clean, stable vocoder", CYAN),
        ("Phonemizer", "Stock en-us (Path A); custom eSpeak optional (Path B)", VIOLET),
        ("Quality tier", "Medium (22 kHz) by default, Low fallback", AMBER),
        ("Deployment", "INT8 ONNX on ONNX Runtime Web, in a Web Worker, cached", SKY),
        ("Validation", "WER < 5%, MCD < 3.0, plus human MOS & A/B tests", PINK),
    ]
    y = Inches(1.95)
    rh = Inches(0.74)
    for name, val, col in rows:
        add_rect(s, Inches(0.55), y, Inches(12.23), rh - Inches(0.1),
                 fill=CARD_BG2, line=col, line_w=Pt(1), radius=0.08)
        add_rect(s, Inches(0.55), y, Inches(0.06), rh - Inches(0.1), fill=col, line=None, radius=0)
        add_txt(s, name, Inches(0.8), y, Inches(3.0), rh - Inches(0.1),
                size=Pt(13.5), color=col, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_txt(s, val, Inches(3.9), y, Inches(8.7), rh - Inches(0.1),
                size=Pt(12.5), color=SL200, anchor=MSO_ANCHOR.MIDDLE)
        y = y + rh
    footer(s, 14)


def slide_15_closing(prs):
    s = new_slide(prs)
    add_txt(s, "Indian English, on-device, done right",
            Inches(1), Inches(2.15), Inches(11.33), Inches(0.9),
            size=Pt(38), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_txt(s, "eSpeak-ng picks the sounds  ·  a fine-tuned VITS makes them Indian  ·  it all runs in the browser",
            Inches(1), Inches(3.15), Inches(11.33), Inches(0.5),
            size=Pt(14), color=CYAN, align=PP_ALIGN.CENTER)

    stats = [("3–5 hrs", "of audio", EMERALD), ("4–6 hrs", "to train", CYAN),
             ("~15 MB", "in browser", AMBER), ("< 150 ms", "to first sound", VIOLET),
             ("100%", "client-side", PINK)]
    cw, gap = Inches(2.05), Inches(0.2)
    total = len(stats) * cw + (len(stats) - 1) * gap
    x0 = (SW - total) // 2
    for i, (big, small, col) in enumerate(stats):
        x = x0 + i * (cw + gap)
        add_rect(s, x, Inches(4.15), cw, Inches(1.3), fill=CARD_BG2, line=col, line_w=Pt(1), radius=0.1)
        add_txt(s, big, x, Inches(4.32), cw, Inches(0.6), size=Pt(22), color=col, bold=True, align=PP_ALIGN.CENTER)
        add_txt(s, small, x, Inches(4.95), cw, Inches(0.4), size=Pt(11), color=SL400, align=PP_ALIGN.CENTER)

    add_txt(s, "Thank you  ·  Questions welcome",
            Inches(1), Inches(6.0), Inches(11.33), Inches(0.5),
            size=Pt(16), color=SL200, bold=True, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_bigidea(prs)
    slide_04_pipeline(prs)
    slide_05_levers(prs)
    slide_06_eninsight(prs)
    slide_07_twopaths(prs)
    slide_08_data(prs)
    slide_09_finetune(prs)
    slide_10_hyperparams(prs)
    slide_11_eval(prs)
    slide_12_deploy(prs)
    slide_13_tradeoffs(prs)
    slide_14_recommendation(prs)
    slide_15_closing(prs)

    prs.save(str(OUTPUT_PATH))
    print(f"Saved: {OUTPUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
